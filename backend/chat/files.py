"""On-demand file generation for chat: an actual file, not a description of one.

Deliberately not persisted anywhere - generated on demand and held in an
in-memory cache with a short TTL, served once. Render's free-tier disk is
ephemeral, so pretending this is a durable store would be dishonest; a chat
turn is not a document management system. Ask again and it regenerates.

Two content shapes cover every format here: tabular (columns/rows - xlsx,
csv) and document (title/sections of heading+paragraphs - docx, pdf, pptx,
md, txt; a pptx slide is just a section with its paragraphs as bullets).
"""

import io
import logging
import time
import uuid
from collections.abc import Callable
from typing import Literal

from backend.handoff import extract_json
from backend.providers import registry

log = logging.getLogger("aco.files")

Kind = Literal["xlsx", "csv", "docx", "pdf", "pptx", "md", "txt"]

_EXT_MARKERS: dict[Kind, tuple[str, ...]] = {
    "xlsx": (".xlsx",), "csv": (".csv",), "docx": (".docx",),
    "pdf": (".pdf",), "pptx": (".pptx",), "md": (".md",), "txt": (".txt",),
}
# Word-phrase markers are gated behind a make-verb (below) since bare words
# like "presentation" or "document" are common outside a file-making request
# ("explain your presentation of the results") - extension mentions above are
# unambiguous enough to skip that gate.
_WORD_MARKERS: dict[Kind, tuple[str, ...]] = {
    "xlsx": ("excel", "spreadsheet", "excel sheet", "excel file"),
    "csv": ("csv file", "csv sheet", "as a csv", "in csv format"),
    "docx": ("word doc", "word document", "docx file"),
    "pdf": ("pdf report", "pdf document", "pdf file", "as a pdf"),
    "pptx": ("powerpoint", "power point", "slide deck", "presentation", "slides for"),
    "md": ("markdown file", "markdown document"),
    "txt": ("text file", "plain text file"),
}
_MAKE_VERBS = (
    "make", "create", "generate", "build", "write", "give me", "produce", "draft", "export",
)


def wants_file(question: str) -> Kind | None:
    """Cheap keyword gate, same pattern as backend/providers/websearch.py's
    needs_fresh_data - most questions are not asking for a file."""
    low = f" {question.lower()} "
    for kind, exts in _EXT_MARKERS.items():
        if any(e in low for e in exts):
            return kind
    if not any(v in low for v in _MAKE_VERBS):
        return None
    for kind, words in _WORD_MARKERS.items():
        if any(w in low for w in words):
            return kind
    return None


TABULAR_SCHEMA_HINT = (
    'Return JSON only, exactly this shape:\n'
    '{"sheet_name": "short name", "columns": ["Column A", "Column B"], '
    '"rows": [["value", 1], ["value", 2]]}\n'
    "Real data only - invent nothing you cannot support from the question. "
    "Each row must have the same number of cells as there are columns."
)
DOCUMENT_SCHEMA_HINT = (
    'Return JSON only, exactly this shape:\n'
    '{"title": "Document title", "sections": [{"heading": "Section", '
    '"paragraphs": ["First paragraph.", "Second paragraph."]}]}\n'
    "Real content only - invent nothing you cannot support from the question. "
    "For a slide deck, each section becomes one slide - keep paragraphs short "
    "(bullet-length, not prose)."
)
_TABULAR_KINDS = {"xlsx", "csv"}


def schema_hint_for(kind: Kind) -> str:
    return TABULAR_SCHEMA_HINT if kind in _TABULAR_KINDS else DOCUMENT_SCHEMA_HINT


def render_xlsx(data: dict) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = str(data.get("sheet_name") or "Sheet1")[:31]
    columns = [str(c) for c in (data.get("columns") or [])]
    if columns:
        ws.append(columns)
    for row in data.get("rows") or []:
        ws.append([c if isinstance(c, int | float) else str(c) for c in row])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def render_csv(data: dict) -> bytes:
    import csv as csv_mod

    buf = io.StringIO()
    writer = csv_mod.writer(buf)
    columns = [str(c) for c in (data.get("columns") or [])]
    if columns:
        writer.writerow(columns)
    for row in data.get("rows") or []:
        writer.writerow(row)
    return buf.getvalue().encode("utf-8")


def render_docx(data: dict) -> bytes:
    from docx import Document

    doc = Document()
    if data.get("title"):
        doc.add_heading(str(data["title"]), level=0)
    for section in data.get("sections") or []:
        if section.get("heading"):
            doc.add_heading(str(section["heading"]), level=1)
        for para in section.get("paragraphs") or []:
            doc.add_paragraph(str(para))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def render_pdf(data: dict) -> bytes:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    # multi_cell leaves the cursor wherever the text ended, not back at the
    # left margin, unless told to - the next call then has less and less
    # horizontal room until fpdf2 raises "not enough horizontal space".
    # new_x/new_y force a reset to the left margin on the next line every time.
    def line(text: str) -> None:
        pdf.multi_cell(0, 8, text, new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf = FPDF()
    pdf.add_page()
    if data.get("title"):
        pdf.set_font("Helvetica", "B", 16)
        line(str(data["title"]))
        pdf.ln(4)
    for section in data.get("sections") or []:
        if section.get("heading"):
            pdf.set_font("Helvetica", "B", 13)
            line(str(section["heading"]))
        pdf.set_font("Helvetica", "", 11)
        for para in section.get("paragraphs") or []:
            line(str(para))
            pdf.ln(2)
    return bytes(pdf.output())


def render_pptx(data: dict) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    if data.get("title"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = str(data["title"])
    for section in data.get("sections") or []:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(section.get("heading") or "")
        paras = [str(p) for p in (section.get("paragraphs") or [])]
        if paras:
            body = slide.placeholders[1].text_frame
            body.text = paras[0]
            for p in paras[1:]:
                body.add_paragraph().text = p
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_text(data: dict, *, markdown: bool) -> bytes:
    lines: list[str] = []
    if data.get("title"):
        lines.append(f"# {data['title']}" if markdown else str(data["title"]))
        lines.append("")
    for section in data.get("sections") or []:
        if section.get("heading"):
            lines.append(f"## {section['heading']}" if markdown else str(section["heading"]))
            lines.append("")
        for para in section.get("paragraphs") or []:
            lines.append(str(para))
            lines.append("")
    return "\n".join(lines).encode("utf-8")


_RENDERERS: dict[Kind, Callable[[dict], bytes]] = {
    "xlsx": render_xlsx,
    "csv": render_csv,
    "docx": render_docx,
    "pdf": render_pdf,
    "pptx": render_pptx,
    "md": lambda d: render_text(d, markdown=True),
    "txt": lambda d: render_text(d, markdown=False),
}


def _safe_filename(name: str, kind: Kind) -> str:
    safe = "".join(ch if ch.isalnum() or ch in " _-" else "_" for ch in name).strip()
    return f"{safe or 'download'}.{kind}"


async def generate_file(question: str, kind: Kind, model: str) -> tuple[bytes, str] | None:
    """One structured-JSON model call, rendered to real bytes. None on any
    failure - the turn falls back to a normal text answer rather than erroring."""
    try:
        c = await registry.complete(
            model,
            [{"role": "user", "content": f"{question}\n\n{schema_hint_for(kind)}"}],
            temperature=0.2, max_tokens=1500, json_mode=True,
        )
        data, err = extract_json(c.text)
        if err or not isinstance(data, dict):
            log.warning("file generation returned unusable JSON (%s)", err)
            return None
        content = _RENDERERS[kind](data)
        name = str(data.get("sheet_name") or data.get("title") or "download")
        return content, _safe_filename(name, kind)
    except Exception as e:
        log.warning("file generation failed (%s), falling back to a text answer", e)
        return None


# ---------- in-memory, short-TTL cache: generate-on-demand, not a document store ----------

_TTL_S = 600
_MAX_ENTRIES = 200
MIME: dict[Kind, str] = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "csv": "text/csv",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "md": "text/markdown",
    "txt": "text/plain",
}
# file_id -> (bytes, filename, mime, expires_at)
_cache: dict[str, tuple[bytes, str, str, float]] = {}


def _evict_expired() -> None:
    now = time.time()
    for k in [k for k, v in _cache.items() if v[3] < now]:
        _cache.pop(k, None)


def store_file(content: bytes, filename: str, kind: Kind) -> str:
    _evict_expired()
    if len(_cache) >= _MAX_ENTRIES:
        oldest = min(_cache, key=lambda k: _cache[k][3])
        _cache.pop(oldest, None)
    file_id = uuid.uuid4().hex
    _cache[file_id] = (content, filename, MIME[kind], time.time() + _TTL_S)
    return file_id


def get_file(file_id: str) -> tuple[bytes, str, str] | None:
    _evict_expired()
    entry = _cache.get(file_id)
    return entry[:3] if entry else None
